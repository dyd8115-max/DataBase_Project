from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

DARK_BLUE = RGBColor(0x1e, 0x3a, 0x5f)
BLUE = RGBColor(0x25, 0x63, 0xeb)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)
DARK_TEXT = RGBColor(0x1e, 0x29, 0x3b)
ACCENT = RGBColor(0x16, 0xa3, 0x4a)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank = prs.slide_layouts[6]  # blank


def add_slide():
    return prs.slides.add_slide(blank)


def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, left, top, width, height, text, font_size=20, bold=False,
        color=DARK_TEXT, bg_color=None, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    if bg_color:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg_color
    return txBox


def rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ── Slide 1: 표지 ──────────────────────────────────────────
s = add_slide()
bg(s, DARK_BLUE)
rect(s, 0, 5.8, 13.33, 1.7, BLUE)
box(s, 1, 1.2, 11, 1.2, "데이터베이스 과제", 28, bold=True, color=RGBColor(0xAD,0xC8,0xFF), align=PP_ALIGN.CENTER)
box(s, 1, 2.4, 11, 1.8, "도서관 관리 시스템", 52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
box(s, 1, 4.3, 11, 0.6, "PostgreSQL + Python Flask 기반 웹 서비스", 22, color=RGBColor(0xCB,0xD5,0xE1), align=PP_ALIGN.CENTER)
box(s, 1, 6.1, 11, 0.5, "제출일: 2026년 6월 10일", 16, color=WHITE, align=PP_ALIGN.CENTER)

# ── Slide 2: 서비스 개요 ───────────────────────────────────
s = add_slide()
bg(s, LIGHT_GRAY)
rect(s, 0, 0, 13.33, 1.2, DARK_BLUE)
box(s, 0.5, 0.25, 12, 0.8, "서비스 개요", 30, bold=True, color=WHITE)

items = [
    ("도서 관리",    "도서 등록·조회·삭제, 카테고리/키워드 검색"),
    ("회원 관리",    "회원 등록·조회, 대출 이력 집계"),
    ("대출 처리",    "트랜잭션 기반 대출 처리 — 재고 자동 차감"),
    ("반납 처리",    "반납 시 재고 복구, 연체 여부 표시"),
    ("대시보드",     "전체 도서·회원·대출 중·연체 건수 통계"),
]
for i, (title, desc) in enumerate(items):
    top = 1.5 + i * 1.0
    rect(s, 0.6, top, 3.5, 0.72, DARK_BLUE)
    box(s, 0.7, top + 0.1, 3.3, 0.55, title, 18, bold=True, color=WHITE)
    box(s, 4.4, top + 0.12, 8.5, 0.55, desc, 17, color=DARK_TEXT)

# ── Slide 3: 기술 스택 ────────────────────────────────────
s = add_slide()
bg(s, LIGHT_GRAY)
rect(s, 0, 0, 13.33, 1.2, DARK_BLUE)
box(s, 0.5, 0.25, 12, 0.8, "기술 스택", 30, bold=True, color=WHITE)

stacks = [
    ("Database",  "PostgreSQL",        "릴레이션, 트랜잭션, 인덱스"),
    ("Backend",   "Python 3 + Flask",  "REST 라우팅, DB 연동 (psycopg2)"),
    ("Frontend",  "Bootstrap 5",       "반응형 UI, Jinja2 템플릿"),
    ("배포 환경", "Linux + Python venv","로컬 실행 / python app.py"),
]
for i, (layer, tech, detail) in enumerate(stacks):
    top = 1.5 + i * 1.25
    rect(s, 0.6, top, 2.2, 0.8, BLUE)
    box(s, 0.7, top + 0.15, 2.0, 0.55, layer, 17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    box(s, 3.1, top + 0.0, 3.8, 0.45, tech, 20, bold=True, color=DARK_BLUE)
    box(s, 3.1, top + 0.42, 9.5, 0.4, detail, 15, color=DARK_TEXT)

# ── Slide 4: DB 설계 — 릴레이션 ──────────────────────────
s = add_slide()
bg(s, LIGHT_GRAY)
rect(s, 0, 0, 13.33, 1.2, DARK_BLUE)
box(s, 0.5, 0.25, 12, 0.8, "데이터베이스 설계 — 릴레이션 (Relation)", 30, bold=True, color=WHITE)

tables = [
    ("members", "회원", ["member_id  PK", "name  NOT NULL", "email  UNIQUE", "phone", "join_date  DEFAULT"]),
    ("books",   "도서", ["book_id  PK", "title / author", "isbn  UNIQUE", "category", "total_copies / avail_copies  CHECK"]),
    ("loans",   "대출", ["loan_id  PK", "book_id  FK → books", "member_id  FK → members", "loan_date / due_date", "status  CHECK (active/returned/overdue)"]),
]
for i, (tname, kor, cols) in enumerate(tables):
    left = 0.4 + i * 4.3
    rect(s, left, 1.35, 4.0, 0.55, BLUE)
    box(s, left + 0.1, 1.38, 3.8, 0.5, f"{tname}  ({kor})", 18, bold=True, color=WHITE)
    for j, col in enumerate(cols):
        rect(s, left, 1.9 + j * 0.85, 4.0, 0.78, WHITE)
        box(s, left + 0.15, 1.98 + j * 0.85, 3.7, 0.6, col, 14, color=DARK_TEXT)

# ── Slide 5: 쿼리 예시 ────────────────────────────────────
s = add_slide()
bg(s, LIGHT_GRAY)
rect(s, 0, 0, 13.33, 1.2, DARK_BLUE)
box(s, 0.5, 0.25, 12, 0.8, "주요 쿼리 (Query)", 30, bold=True, color=WHITE)

queries = [
    ("JOIN 쿼리 — 대출 목록 조회",
     "SELECT l.loan_id, b.title, m.name, l.due_date\nFROM loans l\nJOIN books b ON l.book_id = b.book_id\nJOIN members m ON l.member_id = m.member_id\nORDER BY l.loan_id DESC;"),
    ("GROUP BY 집계 — 회원별 대출 횟수",
     "SELECT m.name,\n  COUNT(l.loan_id) AS total_loans,\n  COUNT(CASE WHEN l.status='active' THEN 1 END) AS active\nFROM members m\nLEFT JOIN loans l ON m.member_id = l.member_id\nGROUP BY m.member_id;"),
]
for i, (title, sql) in enumerate(queries):
    top = 1.4 + i * 2.8
    box(s, 0.6, top, 12, 0.45, title, 17, bold=True, color=DARK_BLUE)
    rect(s, 0.6, top + 0.48, 12.1, 2.1, RGBColor(0x1e,0x29,0x3b))
    box(s, 0.85, top + 0.55, 11.7, 2.0, sql, 14, color=RGBColor(0x7D,0xD3,0xFC))

# ── Slide 6: 트랜잭션 ─────────────────────────────────────
s = add_slide()
bg(s, LIGHT_GRAY)
rect(s, 0, 0, 13.33, 1.2, DARK_BLUE)
box(s, 0.5, 0.25, 12, 0.8, "트랜잭션 (Transaction) — 대출 처리", 30, bold=True, color=WHITE)

steps = [
    ("①", "SELECT ... FOR UPDATE",  "books 테이블 행 잠금 — 동시 대출 방지"),
    ("②", "avail_copies 확인",      "재고 0이면 ROLLBACK → 오류 메시지"),
    ("③", "INSERT INTO loans",      "대출 레코드 생성"),
    ("④", "UPDATE books SET avail_copies - 1", "재고 차감"),
    ("⑤", "COMMIT",                 "모든 단계 성공 시 최종 반영"),
]
for i, (num, op, desc) in enumerate(steps):
    top = 1.45 + i * 1.02
    rect(s, 0.5, top, 0.65, 0.75, DARK_BLUE)
    box(s, 0.52, top + 0.12, 0.6, 0.55, num, 18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, 1.3, top, 4.5, 0.75, BLUE)
    box(s, 1.4, top + 0.12, 4.3, 0.55, op, 16, bold=True, color=WHITE)
    box(s, 6.0, top + 0.15, 7.0, 0.55, desc, 16, color=DARK_TEXT)

box(s, 0.5, 6.7, 12, 0.45,
    "반납도 동일 구조: UPDATE loans(returned) + UPDATE books(avail+1) → COMMIT",
    14, color=BLUE)

# ── Slide 7: 화면 구성 ────────────────────────────────────
s = add_slide()
bg(s, LIGHT_GRAY)
rect(s, 0, 0, 13.33, 1.2, DARK_BLUE)
box(s, 0.5, 0.25, 12, 0.8, "화면 구성", 30, bold=True, color=WHITE)

screens = [
    ("대시보드",   "전체 도서·회원·대출 통계 카드 + 최근 대출 테이블"),
    ("도서 관리",  "도서 목록 테이블, 카테고리/키워드 검색, 등록/삭제"),
    ("회원 관리",  "회원 목록, 총 대출 횟수·현재 대출 집계"),
    ("대출/반납",  "대출 현황 테이블, 연체 행 빨간 강조, 반납 버튼"),
    ("신규 대출",  "재고 있는 도서 선택 + 회원 선택 → 트랜잭션 처리"),
]
for i, (title, desc) in enumerate(screens):
    top = 1.45 + i * 1.02
    rect(s, 0.5, top, 2.8, 0.78, DARK_BLUE)
    box(s, 0.6, top + 0.15, 2.6, 0.55, title, 18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    box(s, 3.6, top + 0.18, 9.5, 0.55, desc, 17, color=DARK_TEXT)

# ── Slide 8: 실행 방법 ────────────────────────────────────
s = add_slide()
bg(s, LIGHT_GRAY)
rect(s, 0, 0, 13.33, 1.2, DARK_BLUE)
box(s, 0.5, 0.25, 12, 0.8, "실행 방법", 30, bold=True, color=WHITE)

steps2 = [
    ("1", "DB 생성",       "psql -U postgres\nCREATE DATABASE library_db;"),
    ("2", "스키마 적용",   "psql -U postgres -d library_db -f schema.sql"),
    ("3", "환경 변수",     "cp .env.example .env  →  DB 접속 정보 입력"),
    ("4", "패키지 설치",   "pip install -r requirements.txt"),
    ("5", "서버 실행",     "python app.py  →  http://localhost:5000 접속"),
]
for i, (num, title, cmd) in enumerate(steps2):
    top = 1.4 + i * 1.1
    rect(s, 0.5, top, 0.55, 0.85, BLUE)
    box(s, 0.53, top+0.18, 0.5, 0.5, num, 20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    box(s, 1.25, top+0.02, 2.5, 0.42, title, 17, bold=True, color=DARK_BLUE)
    rect(s, 1.25, top+0.45, 11.5, 0.38, RGBColor(0x1e,0x29,0x3b))
    box(s, 1.4, top+0.47, 11.2, 0.35, cmd, 13, color=RGBColor(0x86,0xEF,0xAC))

# ── Slide 9: 마무리 ───────────────────────────────────────
s = add_slide()
bg(s, DARK_BLUE)
rect(s, 0, 5.5, 13.33, 2.0, BLUE)
box(s, 1, 1.5, 11, 1.2, "학습 내용 정리", 28, bold=True, color=RGBColor(0xAD,0xC8,0xFF), align=PP_ALIGN.CENTER)

points = [
    "릴레이션 설계  —  테이블 구조, 기본키/외래키, 제약 조건",
    "쿼리 작성      —  JOIN, GROUP BY, CASE WHEN, ILIKE",
    "트랜잭션       —  SELECT FOR UPDATE, COMMIT / ROLLBACK",
    "웹 연동        —  Flask + psycopg2로 DBMS와 백엔드 연결",
]
for i, pt in enumerate(points):
    box(s, 2, 2.9 + i * 0.72, 9.5, 0.6, f"✓  {pt}", 18, color=WHITE)

box(s, 1, 5.7, 11, 0.6, "GitHub: github.com/dyd8115-max/DataBase_Project", 16, color=WHITE, align=PP_ALIGN.CENTER)

prs.save("/home/user/DataBase_Project/도서관_관리_시스템_발표자료.pptx")
print("PPT 생성 완료")
